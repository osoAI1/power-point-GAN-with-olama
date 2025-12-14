# ============================================================
#  presentation_creator_rag_pdf.py
#  - Flask App with:
#       * Local GGUF model via llama-cpp-python
#       * Smart LLM text cleaning (removes “Sure, here is…” etc.)
#       * Optional PDF upload as knowledge source
#       * TF-IDF based chunk retrieval (RAG) per slide
#       * Temp files auto cleanup
# ============================================================

import os
import sys
import threading
import tempfile
import atexit
from wsgiref import simple_server
import webbrowser
import time
import re

from flask import Flask, request, send_file, jsonify, render_template
from llama_cpp import Llama
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from PyPDF2 import PdfReader            # pip install PyPDF2
from sklearn.feature_extraction.text import TfidfVectorizer  # pip install scikit-learn
import numpy as np

# ==================== CONFIGURATION ====================
MODEL_PATH = "qwen2.gguf"  # اسم ملف الموديل GGUF
MODEL_LOADED = False
TEMP_FILES = []  # Track temp files for cleanup


def cleanup_temp_files():
    """Clean up temporary files on exit"""
    for file_path in list(TEMP_FILES):
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"🗑️  Cleaned up: {file_path}")
                TEMP_FILES.remove(file_path)
        except Exception as e:
            print(f"⚠️  Failed to cleanup {file_path}: {e}")


# Register cleanup function
atexit.register(cleanup_temp_files)


# ==================== TEXT CLEANING ====================
def clean_text(text):
    if not isinstance(text, str):
        return text

    # عبارات LLM الشائعة اللي عايز تشيلها
    bad_patterns = [
        r"Sure,? here.?s .*?:?",
        r"Certainly!?.*?:?",
        r"Certainly,.*?:?",
        r"Here.? is .*?:?",
        r"Here.?s .*?:?",
        r"As an AI language model.*",
        r"Make sure to .*",
        r"End with an intriguing question.*",
        r"Use appropriate vocabulary.*",
        r"Highlight key concepts.*",
        r"Below is .*",
        r"This slide covers .*",
        r"Here is some content.*",
        r"Here is the content.*",
        r"Here.?s an academic explanation.*",
        r"Sure!?.*?:?",
        r"Sure,.*",
        r"Do not include citations.*",
        r"Do not include references.*",
        r"Act as .*?:?",
        r"Provide an overview.*",
        r"Make sure examples are relevant.*",
        r"Write content for the slide.*",
        r"Explain the topic.*",
        r"Use a professional writing style.*",
        r"Use formal academic tone.*"
    ]

    for pat in bad_patterns:
        text = re.sub(pat, "", text, flags=re.IGNORECASE)

    # إزالة الرموز المزعجة
    text = re.sub(r"[{}#•\-*\"']", "", text)

    # إزالة [1] [23] وغيرها
    text = re.sub(r"\[\s*\d+\s*\]", "", text)

    # إزالة المسافات الزائدة
    text = re.sub(r"\s+", " ", text).strip()

    return text


# ==================== PDF HANDLING ====================
def extract_text_from_pdf(file_path, max_chars=40000):
    """Extract plain text from PDF."""
    try:
        reader = PdfReader(file_path)
        texts = []
        total = 0
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                texts.append(t)
                total += len(t)
                if total >= max_chars:
                    break
        full_text = "\n".join(texts)
        print(f"📄 PDF extracted chars: {len(full_text)}")
        return full_text[:max_chars]
    except Exception as e:
        print(f"⚠️ Failed to read PDF: {e}")
        return ""


# ==================== SLIDE SPLITTING ====================
MAX_SLIDE_CHARS = 1000


def split_content_if_needed(text):
    text = clean_text(text)
    if len(text) <= MAX_SLIDE_CHARS:
        return [text]

    # Split on sentence boundaries
    sentences = re.split(r'(?<=[.!?؟]) +', text)
    slides = []
    current = ""

    for s in sentences:
        if len(current) + len(s) <= MAX_SLIDE_CHARS:
            if current:
                current += " " + s
            else:
                current = s
        else:
            slides.append(current.strip())
            current = s

    if current.strip():
        slides.append(current.strip())

    return slides


# ==================== AI CONTENT GENERATOR ====================
class ContentGenerator:
    def __init__(self, model_path):
        if not os.path.exists(model_path):
            raise FileNotFoundError(f"❌ Model not found: {model_path}")
        print("🧠 Loading AI model... (This may take a minute)")
        try:
            self.llm = Llama(
                model_path=model_path,
                n_ctx=4096,
                n_gpu_layers=32,
                temperature=0.25,
                top_p=0.9,
                verbose=False,
                # repeat_penalty=1.25
            )
            # PDF index (RAG)
            self.pdf_chunks = []
            self.pdf_vectorizer = None
            self.pdf_tfidf = None

            global MODEL_LOADED
            MODEL_LOADED = True
            print("✅ Model loaded successfully!")
        except Exception as e:
            print(f"❌ Failed to load model: {e}")
            raise

    # -------- PDF RAG INDEXING --------
    def clear_pdf_index(self):
        self.pdf_chunks = []
        self.pdf_vectorizer = None
        self.pdf_tfidf = None
        print("📚 PDF index cleared (no PDF / not used).")

    def build_pdf_index(self, text: str, chunk_size: int = 550, overlap: int = 120):
        """Turn raw pdf text into overlapping chunks + TF-IDF index."""
        if not text or not text.strip():
            self.clear_pdf_index()
            return

        # Normalize spaces
        plain = re.sub(r"\s+", " ", text).strip()
        chunks = []
        start = 0
        n = len(plain)

        while start < n:
            end = min(start + chunk_size, n)
            chunk = plain[start:end].strip()
            if len(chunk) > 50:  # تجاهل الشظايا الصغيرة جدًا
                chunks.append(chunk)
            if end == n:
                break
            start = end - overlap  # overlap بسيط بينهم

        if not chunks:
            self.clear_pdf_index()
            return

        print(f"📚 Building TF-IDF index for {len(chunks)} PDF chunks...")
        try:
            vectorizer = TfidfVectorizer()
            tfidf = vectorizer.fit_transform(chunks)
            self.pdf_chunks = chunks
            self.pdf_vectorizer = vectorizer
            self.pdf_tfidf = tfidf
            print("✅ PDF TF-IDF index built successfully.")
        except Exception as e:
            print(f"⚠️ Failed to build TF-IDF index: {e}")
            self.clear_pdf_index()

    def set_pdf_from_text(self, text: str):
        """Public API: takes raw text from PDF, builds TF-IDF index."""
        if text and text.strip():
            self.build_pdf_index(text)
        else:
            self.clear_pdf_index()

    def _build_reference_snippet(self, topic: str, slide_title: str, max_ref_chars: int = 1400) -> str:
        """Pick most relevant chunks for this slide using TF-IDF cosine similarity."""
        if self.pdf_vectorizer is None or self.pdf_tfidf is None or not self.pdf_chunks:
            return ""

        query = f"""Topic: {topic} Slide focus: {slide_title} 
        Keywords to match: core concepts, definitions, examples, applications."""

        try:
            q_vec = self.pdf_vectorizer.transform([query])
            # cosine similarity for normalized tf-idf is just dot product
            sims = (self.pdf_tfidf @ q_vec.T).toarray().ravel()
        except Exception as e:
            print(f"⚠️ Failed to compute TF-IDF similarity: {e}")
            return ""

        # Sort chunks by similarity descending
        indices = np.argsort(-sims)
        selected = []
        total_chars = 0

        for idx in indices:
            if sims[idx] <= 0:
                # no more useful matches
                break
            ch = self.pdf_chunks[idx]
            selected.append(ch)
            total_chars += len(ch)
            if total_chars >= max_ref_chars:
                break

        if not selected:
            # fallback: first few chunks
            joined = " ".join(self.pdf_chunks[:4])
            return joined[:max_ref_chars]

        ref_text = " ".join(selected)
        return ref_text[:max_ref_chars]

    # -------- MAIN GENERATION --------
    def generate_slide_content(self, topic, slide_title, prompt_style):
        # ===== PDF reference context (لو موجود وتم بناؤه) =====
        reference_block = ""
        if self.pdf_vectorizer is not None and self.pdf_tfidf is not None:
            ref_snippet = self._build_reference_snippet(topic, slide_title)
            if ref_snippet:
                reference_block = f"""
You have the following reference material extracted from a PDF uploaded by the lecturer.
Use it as the MAIN knowledge source and stay as close as possible to its concepts,
definitions, and terminology. Do NOT mention the PDF itself or say that you are using
reference material. Just write clean academic content.

REFERENCE MATERIAL:
\"\"\"{ref_snippet}\"\"\"


"""

        # 🔥 النسخة الذكية التي تمنع أي كلام LLM من الأساس
        prompt1 = f"""{reference_block}You are a university lecturer.

Write ONLY the final academic content for the slide titled: '{slide_title}'
in a lecture on: '{topic}'.

STRICT RULES:
- DO NOT start with phrases like: "Sure", "Certainly", "Here is", "Make sure", 
  "In this slide", "This content", "Below is", "The following".
- DO NOT explain what you are doing.
- DO NOT talk about generating content.
- DO NOT give instructions or meta-comments.
- DO NOT introduce or summarize the slide meta.
- DO NOT include helper text, transitions, or disclaimers.
- ONLY write clean academic paragraphs.

CONTENT REQUIREMENTS:
- 2–3 paragraphs.
- Each paragraph must contain 4–6 sentences.
- Use formal academic tone.
- Provide explanation + example + real application.
- No bullet points or lists.

Begin immediately with the academic content:
"""

        prompt2 = f"""{reference_block}Act only as a domain expert writing polished academic content.

Write ONLY the final paragraphs for the slide titled '{slide_title}' 
in a lecture about '{topic}'.

Forbidden:
- Any AI disclaimers.
- Any helper phrases ("Here is", "Sure", "Certainly").
- Any meta-instructions.
- Any slide descriptions.
- Any notes about what you will do.
- Any teaching/explanation about the slide structure.

Required:
- 2–3 academic paragraphs (4–6 sentences each).
- High-density, expert-level writing.
- Clear reasoning + example + real-world application.
- No lists, bullets, or headings.

Start directly with the paragraph content:
"""

        prompt = prompt1 if prompt_style == "prompt1" else prompt2

        try:
            print(f"   🤖 Generating content for: {slide_title}")
            out = self.llm(prompt, max_tokens=600)

            text = clean_text(out["choices"][0]["text"].strip())
            paragraphs = [p.strip() for p in text.split("\n") if len(p.strip()) > 40]

            return paragraphs[:3] if paragraphs else ["Content generation failed."]
        except Exception as e:
            print(f"   ❌ Error generating content: {e}")
            return [f"Error generating content: {str(e)}"]


# ==================== PPT CREATOR ====================

# ==================== PPT CREATOR (شكل ثابت واحترافي محسّن) ====================
class PresentationCreator:
    def __init__(self, content_generator):
        self.content_generator = content_generator
        
        # نظام الألوان الثابت الاحترافي (درجات الأزرق)
        self.colors = {
            "title": RGBColor(255, 255, 255),          # أبيض
            "subtitle": RGBColor(173, 216, 230),       # أزرق فاتح
            "text": RGBColor(255, 255, 255),           # أبيض
            "accent": RGBColor(135, 206, 250),         # أزرق سماوي
            "shadow": RGBColor(0, 51, 102)             # أزرق داكن للظلال
        }
        # هنخزن الـ presentation هنا عشان نستخدمه في الشرائح المستمرة
        self.prs = None

    # ======================================================
    #   🔥 تقسيم النص إلى أجزاء لا تتجاوز 850 حرف
    # ======================================================
    def split_into_chunks(self, text, max_chars=700):
        """
        تقسيم النص إلى أجزاء لا تتجاوز max_chars حرف
        """
        chunks = []
        text = text.strip()

        while len(text) > max_chars:
            # محاولة التقسيم عند أقرب نقطة نهاية جملة
            split_pos = text.rfind(".", 0, max_chars)

            if split_pos == -1:
                split_pos = max_chars   # لو مفيش نقطة، نقسم تقسيم مباشر

            chunks.append(text[:split_pos].strip())
            text = text[split_pos:].strip()

        if text:
            chunks.append(text)

        return chunks

    # ======================================================
    #                 إنشاء العرض التقديمي
    # ======================================================
    def create_presentation(self, title, num_slides, prompt_style):
        """
        إنشاء عرض تقديمي ثابت واحترافي
        """
        print(f"🎨 Creating professional presentation: '{title}' ({num_slides} slides)")
        
        # استخدام القالب الموجود
        try:
            prs = Presentation("Blue and White Gradient Artificial Intelligence Presentation.pptx")
        except:
            print("⚠️ Template not found, creating new presentation with default style")
            prs = Presentation()
            
        # نخزن الـ prs في الكائن عشان نستخدمه في الشرائح المستمرة
        self.prs = prs

        # إعداد أبعاد الشرائح (16:9 معيار احترافي)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # ==================== شريحة العنوان ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # إضافة طبقة خلفية زرقاء داكنة لشريحة العنوان
        background = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors["shadow"]
        background.line.fill.background()
        
        # العنوان الرئيسي
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = tb.text_frame
        tf.text = title.upper()
        
        p = tf.paragraphs[0]
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = self.colors["title"]
        p.alignment = PP_ALIGN.CENTER
        
        # إضافة ظل للنص
        try:
            p.font.shadow = True
        except:
            pass
        
        # العنوان الفرعي
        sub_tb = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
        sub_tf = sub_tb.text_frame
        sub_tf.text = "Professional AI-Generated Presentation"
        sub_p = sub_tf.paragraphs[0]
        sub_p.font.size = Pt(22)
        sub_p.font.color.rgb = self.colors["subtitle"]
        sub_p.alignment = PP_ALIGN.CENTER
        sub_p.font.italic = True
        
        # ==================== هيكل الشرائح الذكي ====================
        structure = self._get_slide_structure(num_slides)
        
        # ==================== إنشاء الشرائح المحتوية ====================
        for i, slide_title in enumerate(structure):
            print(f"   📊 Slide {i+2}/{num_slides}: {slide_title}")
            
            # إنشاء شريحة جديدة
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # إضافة خلفية زرقاء داكنة لكل شريحة
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.colors["shadow"]
            bg.line.fill.background()
            
            # إضافة تذييل ثابت
            self._add_footer(slide, title, i+2, num_slides)
            
            # ==================== العنوان الرئيسي للشريحة ====================
            title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
            title_tf = title_tb.text_frame
            title_tf.text = slide_title.upper()
            
            title_p = title_tf.paragraphs[0]
            title_p.font.size = Pt(32)
            title_p.font.bold = True
            title_p.font.color.rgb = self.colors["title"]
            
            # ==================== خط فاصل تحت العنوان ====================
            line = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0.5), Inches(1.4), Inches(12), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.colors["accent"]
            line.line.fill.background()
            
            # ==================== المحتوى من الذكاء الاصطناعي ====================
            paragraphs = self.content_generator.generate_slide_content(
                title, slide_title, prompt_style
            )
            
            all_chunks = []
            if paragraphs:
                for para in paragraphs:
                    # هنا بنقسّم كل فقرة إلى أجزاء لا تزيد عن 850 حرف
                    all_chunks.extend(self.split_into_chunks(para, max_chars=850))
                
            # عرض المحتوى مع تنسيق احترافي
            self._add_formatted_content(slide, all_chunks, slide_title)
        
        # ==================== شريحة النهاية ====================
        if num_slides > 3:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # خلفية
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.colors["shadow"]
            bg.line.fill.background()
            
            # تذييل
            self._add_footer(slide, title, num_slides + 1, num_slides)
            
            # عنوان الشكر
            thanks_tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2))
            thanks_tf = thanks_tb.text_frame
            thanks_tf.text = "THANK YOU"
            
            thanks_p = thanks_tf.paragraphs[0]
            thanks_p.font.size = Pt(64)
            thanks_p.font.bold = True
            thanks_p.font.color.rgb = self.colors["accent"]
            thanks_p.alignment = PP_ALIGN.CENTER
            
            # نص إضافي
            details_tb = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))
            details_tf = details_tb.text_frame
            details_tf.text = "Questions & Discussion\nAI-Generated Content"
            
            details_p = details_tf.paragraphs[0]
            details_p.font.size = Pt(24)
            details_p.font.color.rgb = self.colors["subtitle"]
            details_p.alignment = PP_ALIGN.CENTER
        
        # ==================== حفظ الملف ====================
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        filename = temp_file.name
        temp_file.close()
        prs.save(filename)
        
        TEMP_FILES.append(filename)
        print(f"✅ Professional presentation saved: {filename}")
        print(f"   📊 Total slides created: {len(prs.slides)}")
        
        return filename
    
    def _get_slide_structure(self, num_slides):
        """
        إنشاء هيكل الشرائح حسب العدد
        """
        base_structures = {
            3: ["Introduction", "Main Content", "Conclusion"],
            4: ["Introduction", "Core Concepts", "Applications", "Conclusion"],
            5: ["Introduction", "Background", "Methodology", "Results", "Conclusion"],
            6: ["Introduction", "Problem Statement", "Theoretical Framework", 
                "Analysis", "Findings", "Recommendations"],
            7: ["Introduction", "Literature Review", "Research Questions", 
                "Methodology", "Data Analysis", "Results", "Conclusion"],
            8: ["Introduction", "Background", "Problem Definition", "Objectives",
                "Methodology", "Analysis", "Results", "Conclusion"],
            9: ["Executive Summary", "Introduction", "Market Analysis", 
                "Solution Overview", "Technical Details", "Implementation",
                "Benefits", "Case Studies", "Next Steps"],
            10: ["Title", "Agenda", "Introduction", "Current State",
                 "Challenges", "Proposed Solution", "Technical Architecture",
                 "Implementation Plan", "ROI Analysis", "Conclusion"],
        }
        
        if num_slides in base_structures:
            return base_structures[num_slides][1:]  # نزيل "Title" إذا موجود
        
        # هيكل افتراضي للعدد الأكبر
        return [
            "Introduction",
            "Background & Context",
            "Core Concepts",
            "Methodology",
            "Key Findings",
            "Analysis",
            "Applications",
            "Case Studies",
            "Best Practices",
            "Conclusion & Recommendations"
        ][:num_slides - 1]
    
    def _add_footer(self, slide, title, slide_number, total_slides):
        """
        إضافة تذييل احترافي للشريحة
        """
        # إضافة خط فاصل في الأسفل
        footer_line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(6.8), Inches(12), Inches(0.02)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.colors["accent"]
        footer_line.line.fill.background()
        
        # إضافة رقم الشريحة في الزاوية اليمنى
        num_tb = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1), Inches(0.4))
        num_tf = num_tb.text_frame
        num_tf.text = f"{slide_number}"
        num_p = num_tf.paragraphs[0]
        num_p.font.size = Pt(16)
        num_p.font.bold = True
        num_p.font.color.rgb = self.colors["accent"]
        num_p.alignment = PP_ALIGN.RIGHT
        
        # إضافة العدد الكلي للشرائح
        total_tb = slide.shapes.add_textbox(Inches(11.5), Inches(7.2), Inches(1), Inches(0.4))
        total_tf = total_tb.text_frame
        total_tf.text = f"/ {total_slides}"
        total_p = total_tf.paragraphs[0]
        total_p.font.size = Pt(12)
        total_p.font.color.rgb = self.colors["subtitle"]
        total_p.alignment = PP_ALIGN.RIGHT
        
        # إضافة عنوان العرض في الزاوية اليسرى
        title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(8), Inches(0.4))
        title_tf = title_tb.text_frame
        title_tf.text = title[:40] + ("..." if len(title) > 40 else "")
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(14)
        title_p.font.color.rgb = self.colors["subtitle"]
        
        # إضافة تاريخ اليوم
        import datetime
        date_tb = slide.shapes.add_textbox(Inches(8.5), Inches(7), Inches(3), Inches(0.4))
        date_tf = date_tb.text_frame
        date_tf.text = datetime.datetime.now().strftime("%Y-%m-%d")
        date_p = date_tf.paragraphs[0]
        date_p.font.size = Pt(12)
        date_p.font.color.rgb = self.colors["subtitle"]
        date_p.alignment = PP_ALIGN.CENTER
    
    def _add_formatted_content(self, slide, content_chunks, slide_title):
        """
        إضافة محتوى منسق بشكل احترافي
        """
        if not content_chunks:
            return
        
        content_y = Inches(1.8)  # بداية المحتوى بعد العنوان والخط الفاصل
        
        for chunk_index, chunk in enumerate(content_chunks):
            if chunk_index > 0:
                # شريحة جديدة للمحتوى المستمر
                slide = self._create_continued_slide(slide_title)
                content_y = Inches(1.8)
            
            # إضافة مربع نص للمحتوى
            content_tb = slide.shapes.add_textbox(Inches(0.8), content_y, Inches(11.5), Inches(5))
            content_tf = content_tb.text_frame
            content_tf.word_wrap = True
            content_tf.text = chunk
            
            # تنسيق النص
            content_p = content_tf.paragraphs[0]
            content_p.font.size = Pt(22)
            content_p.font.color.rgb = self.colors["text"]
            content_p.line_spacing = 1.5
            content_p.space_after = Pt(12)
            
            # جعل الفقرة الأولى بخط عريض
            if chunk_index == 0:
                content_p.font.bold = True
            
            # إضافة نقطة زرقاء في بداية الفقرة الأولى
            if chunk_index == 0 and len(slide_title) > 0:
                bullet = slide.shapes.add_shape(
                    2,  # MSO_SHAPE.OVAL
                    Inches(0.6), content_y + Inches(0.1), Inches(0.15), Inches(0.15)
                )
                bullet.fill.solid()
                bullet.fill.fore_color.rgb = self.colors["accent"]
                bullet.line.fill.background()
    
    def _create_continued_slide(self, slide_title):
        """
        إنشاء شريحة مستمرة للمحتوى الطويل
        """
        # نستخدم نفس الـ prs المخزّن في self.prs
        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # إضافة خلفية
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors["shadow"]
        bg.line.fill.background()
        
        # إضافة تذييل بسيط (بدون أرقام دقيقة)
        footer_line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(6.8), Inches(12), Inches(0.02)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.colors["accent"]
        footer_line.line.fill.background()
        
        # عنوان الشريحة المستمرة
        title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        title_tf = title_tb.text_frame
        title_tf.text = f"{slide_title.upper()} (Continued)"
        
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors["title"]
        
        # خط فاصل
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.4), Inches(12), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors["accent"]
        line.line.fill.background()
        
        return slide

# ==================== FLASK APP ====================
def create_app():
    app = Flask(__name__)
    app.config['SECRET_KEY'] = 'ai-presentation-creator-secret-key'
    app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32MB max upload

    # Components are created lazily
    content_gen = ContentGenerator(MODEL_PATH)
    ppt_creator = PresentationCreator(content_gen)

    def initialize_ai_components():
        """Initialize AI components on demand"""
        nonlocal content_gen, ppt_creator
        try:
            print("🚀 Initializing AI components...")
            content_gen = ContentGenerator(MODEL_PATH)
            ppt_creator = PresentationCreator(content_gen)
            return True
        except Exception as e:
            print(f"❌ Failed to initialize AI components: {e}")
            return False

    @app.after_request
    def cleanup_temp_files_after_request(response):
        try:
            for file_path in list(TEMP_FILES):
                if os.path.exists(file_path):
                    os.remove(file_path)
                    TEMP_FILES.remove(file_path)
                    print(f"🗑️  Cleaned up temp file (after_request): {file_path}")
        except Exception as e:
            print(f"⚠️  Error in after_request cleanup: {e}")
        return response

    @app.route('/')
    def index():
        return render_template("ind.html")

    @app.route('/generate', methods=['POST'])
    def generate():
        try:
            nonlocal content_gen, ppt_creator
            if content_gen is None or ppt_creator is None:
                print("🔄 Initializing AI model on first request...")
                if not initialize_ai_components():
                    return jsonify({'error': 'Failed to initialize AI model. Please check if model file exists.'}), 500

            title = (request.form.get('title') or '').strip()
            slides_str = request.form.get('slides', '8')
            prompt_style = request.form.get('prompt_style', 'prompt1')
            use_pdf_flag = request.form.get('use_pdf', '0') == '1'
            pdf_file = request.files.get('pdf_file')
            try:
                slides = int(slides_str)
            except ValueError:
                return jsonify({'error': 'Slides must be an integer.'}), 400

            if not title or len(title) < 3:
                return jsonify({'error': 'Please enter a valid title (minimum 3 characters)'}), 400

            if slides < 3 or slides > 12:
                return jsonify({'error': 'Number of slides must be between 3 and 12'}), 400

            # PDF (اختياري + اختياري الاستخدام)
            
            if use_pdf_flag and pdf_file and pdf_file.filename:
                print("📚 Using uploaded PDF for RAG")
                tmp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
                pdf_path = tmp_pdf.name
                pdf_file.save(pdf_path)
                tmp_pdf.close()
                TEMP_FILES.append(pdf_path)

                pdf_text = extract_text_from_pdf(pdf_path)
                content_gen.set_pdf_from_text(pdf_text)
            else:
                content_gen.clear_pdf_index()
                print("📚 PDF index cleared (no PDF / not used)")


            print(f"🎯 Starting generation: '{title}' ({slides} slides, {prompt_style}, use_pdf={use_pdf_flag})")
            start_time = time.time()

            file_path = ppt_creator.create_presentation(title, slides, prompt_style)

            end_time = time.time()
            print(f"⏱️  Generation completed in {end_time - start_time:.2f} seconds")

            return send_file(
                file_path,
                as_attachment=True,
                download_name=f"{title.replace(' ', '_')}.pptx",
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        except FileNotFoundError:
            return jsonify({'error': f'Model file not found: {MODEL_PATH}. Please download a model and update MODEL_PATH.'}), 500
        except ValueError as e:
            return jsonify({'error': f'Invalid input: {str(e)}'}), 400
        except Exception as e:
            print(f"❌ Unexpected error during generation: {e}")
            return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500

    @app.route('/health')
    def health():
        status = {
            'status': 'ok',
            'model_loaded': content_gen is not None,
            'model_path': MODEL_PATH,
            'model_exists': os.path.exists(MODEL_PATH),
            'temp_files': len(TEMP_FILES)
        }
        return jsonify(status)

    @app.route('/cleanup', methods=['POST'])
    def manual_cleanup():
        """Manually cleanup temp files"""
        count = len(TEMP_FILES)
        cleanup_temp_files()
        return jsonify({'message': f'Cleaned up {count} temporary files'})

    return app


def main():
    print("=" * 60)
    print("🤖 AI PRESENTATION CREATOR (with PDF RAG TF-IDF)")
    print("=" * 60)
    print(f"📂 Model path: {MODEL_PATH}")
    print(f"📁 Model exists: {'✅' if os.path.exists(MODEL_PATH) else '❌'}")

    if not os.path.exists(MODEL_PATH):
        print("\n⚠️  WARNING: Model file not found!")
        print("Please download a GGUF model and:")
        print("1. Place it in the current directory")
        print("2. Update MODEL_PATH in the code")
        print("\nRecommended models:")
        print("- qwen2.5-1.5b-instruct-q4_K_M.gguf (fast)")
        print("- llama-2-7b-chat.Q4_K_M.gguf (balanced)")
        print("- mistral-7b-instruct-v0.2.Q4_K_M.gguf (good quality)")
        print("\nYou can download models from:")
        print("- https://huggingface.co/TheBloke")
        print("- https://huggingface.co/models?sort=trending&search=gguf")

    app = create_app()

    def open_browser():
        time.sleep(1.5)
        try:
            webbrowser.open("http://127.0.0.1:5000")
        except Exception:
            print("⚠️  Could not open browser automatically")

    threading.Thread(target=open_browser, daemon=True).start()

    print("\n" + "=" * 60)
    print("🚀 Starting server...")
    print("🌐 Open your browser to: http://127.0.0.1:5000")
    print("⏳ First request will load the AI model (may take 30-60 seconds)")
    print("🛑 Press Ctrl+C to stop the server")
    print("=" * 60)

    try:
        with simple_server.make_server("127.0.0.1", 5000, app) as httpd:
            httpd.serve_forever()
    except KeyboardInterrupt:
        print("\n\n👋 Server stopped by user")
        print("🧹 Cleaning up temporary files...")
        cleanup_temp_files()
        print("✅ Cleanup complete. Goodbye!")
    except Exception as e:
        print(f"\n❌ Server error: {e}")
        cleanup_temp_files()


if __name__ == "__main__":
    main()

