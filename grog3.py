# ============================================================
#  presentation_creator_rag_pdf.py
#  - Flask App with:
#       * Local GGUF model via llama-cpp-python
#       * Smart LLM text cleaning (removes “Sure, here is…” etc.)
#       * Optional PDF upload as knowledge source
#       * TF-IDF based chunk retrieval (RAG) per slide
#       * Temp files auto cleanup
# ============================================================

import os
import sys
import threading
import tempfile
import atexit
from wsgiref import simple_server
import webbrowser
import time
import re

from flask import Flask, request, send_file, jsonify, render_template_string
from llama_cpp import Llama
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from PyPDF2 import PdfReader            # pip install PyPDF2
from sklearn.feature_extraction.text import TfidfVectorizer  # pip install scikit-learn
import numpy as np

# ==================== CONFIGURATION ====================
MODEL_PATH = "qwen2.gguf"  # اسم ملف الموديل GGUF
MODEL_LOADED = False
TEMP_FILES = []  # Track temp files for cleanup


def cleanup_temp_files():
    """Clean up temporary files on exit"""
    for file_path in list(TEMP_FILES):
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"🗑️  Cleaned up: {file_path}")
                TEMP_FILES.remove(file_path)
        except Exception as e:
            print(f"⚠️  Failed to cleanup {file_path}: {e}")


# Register cleanup function
atexit.register(cleanup_temp_files)


# ==================== TEXT CLEANING ====================
def clean_text(text):
    if not isinstance(text, str):
        return text

    # عبارات LLM الشائعة اللي عايز تشيلها
    bad_patterns = [
        r"Sure,? here.?s .*?:?",
        r"Certainly!?.*?:?",
        r"Certainly,.*?:?",
        r"Here.? is .*?:?",
        r"Here.?s .*?:?",
        r"As an AI language model.*",
        r"Make sure to .*",
        r"End with an intriguing question.*",
        r"Use appropriate vocabulary.*",
        r"Highlight key concepts.*",
        r"Below is .*",
        r"This slide covers .*",
        r"Here is some content.*",
        r"Here is the content.*",
        r"Here.?s an academic explanation.*",
        r"Sure!?.*?:?",
        r"Sure,.*",
        r"Do not include citations.*",
        r"Do not include references.*",
        r"Act as .*?:?",
        r"Provide an overview.*",
        r"Make sure examples are relevant.*",
        r"Write content for the slide.*",
        r"Explain the topic.*",
        r"Use a professional writing style.*",
        r"Use formal academic tone.*"
    ]

    for pat in bad_patterns:
        text = re.sub(pat, "", text, flags=re.IGNORECASE)

    # إزالة الرموز المزعجة
    text = re.sub(r"[{}#•\-*\"']", "", text)

    # إزالة [1] [23] وغيرها
    text = re.sub(r"\[\s*\d+\s*\]", "", text)

    # إزالة المسافات الزائدة
    text = re.sub(r"\s+", " ", text).strip()

    return text


# ==================== PDF HANDLING ====================
def extract_text_from_pdf(file_path, max_chars=40000):
    """Extract plain text from PDF."""
    try:
        reader = PdfReader(file_path)
        texts = []
        total = 0
        for page in reader.pages:
            t = page.extract_text() or ""
            if t.strip():
                texts.append(t)
                total += len(t)
                if total >= max_chars:
                    break
        full_text = "\n".join(texts)
        print(f"📄 PDF extracted chars: {len(full_text)}")
        return full_text[:max_chars]
    except Exception as e:
        print(f"⚠️ Failed to read PDF: {e}")
        return ""


# ==================== SLIDE SPLITTING ====================
MAX_SLIDE_CHARS = 1300


def split_content_if_needed(text):
    text = clean_text(text)
    if len(text) <= MAX_SLIDE_CHARS:
        return [text]

    # Split on sentence boundaries
    sentences = re.split(r'(?<=[.!?؟]) +', text)
    slides = []
    current = ""

    for s in sentences:
        if len(current) + len(s) <= MAX_SLIDE_CHARS:
            if current:
                current += " " + s
            else:
                current = s
        else:
            slides.append(current.strip())
            current = s

    if current.strip():
        slides.append(current.strip())

    return slides


# ==================== AI CONTENT GENERATOR ====================
class ContentGenerator:
    def __init__(self, model_path):
        if not os.path.exists(model_path):
            raise FileNotFoundError(f"❌ Model not found: {model_path}")
        print("🧠 Loading AI model... (This may take a minute)")
        try:
            self.llm = Llama(
                model_path=model_path,
                n_ctx=4096,
                n_gpu_layers=0,
                temperature=0.25,
                top_p=0.9,
                verbose=False
            )
            # PDF index (RAG)
            self.pdf_chunks = []
            self.pdf_vectorizer = None
            self.pdf_tfidf = None

            global MODEL_LOADED
            MODEL_LOADED = True
            print("✅ Model loaded successfully!")
        except Exception as e:
            print(f"❌ Failed to load model: {e}")
            raise

    # -------- PDF RAG INDEXING --------
    def clear_pdf_index(self):
        self.pdf_chunks = []
        self.pdf_vectorizer = None
        self.pdf_tfidf = None
        print("📚 PDF index cleared (no PDF / not used).")

    def build_pdf_index(self, text: str, chunk_size: int = 600, overlap: int = 120):
        """Turn raw pdf text into overlapping chunks + TF-IDF index."""
        if not text or not text.strip():
            self.clear_pdf_index()
            return

        # Normalize spaces
        plain = re.sub(r"\s+", " ", text).strip()
        chunks = []
        start = 0
        n = len(plain)

        while start < n:
            end = min(start + chunk_size, n)
            chunk = plain[start:end].strip()
            if len(chunk) > 50:  # تجاهل الشظايا الصغيرة جدًا
                chunks.append(chunk)
            if end == n:
                break
            start = end - overlap  # overlap بسيط بينهم

        if not chunks:
            self.clear_pdf_index()
            return

        print(f"📚 Building TF-IDF index for {len(chunks)} PDF chunks...")
        try:
            vectorizer = TfidfVectorizer()
            tfidf = vectorizer.fit_transform(chunks)
            self.pdf_chunks = chunks
            self.pdf_vectorizer = vectorizer
            self.pdf_tfidf = tfidf
            print("✅ PDF TF-IDF index built successfully.")
        except Exception as e:
            print(f"⚠️ Failed to build TF-IDF index: {e}")
            self.clear_pdf_index()

    def set_pdf_from_text(self, text: str):
        """Public API: takes raw text from PDF, builds TF-IDF index."""
        if text and text.strip():
            self.build_pdf_index(text)
        else:
            self.clear_pdf_index()

    def _build_reference_snippet(self, topic: str, slide_title: str, max_ref_chars: int = 1400) -> str:
        """Pick most relevant chunks for this slide using TF-IDF cosine similarity."""
        if self.pdf_vectorizer is None or self.pdf_tfidf is None or not self.pdf_chunks:
            return ""

        query = f"{topic}. {slide_title}"
        try:
            q_vec = self.pdf_vectorizer.transform([query])
            # cosine similarity for normalized tf-idf is just dot product
            sims = (self.pdf_tfidf @ q_vec.T).toarray().ravel()
        except Exception as e:
            print(f"⚠️ Failed to compute TF-IDF similarity: {e}")
            return ""

        # Sort chunks by similarity descending
        indices = np.argsort(-sims)
        selected = []
        total_chars = 0

        for idx in indices:
            if sims[idx] <= 0:
                # no more useful matches
                break
            ch = self.pdf_chunks[idx]
            selected.append(ch)
            total_chars += len(ch)
            if total_chars >= max_ref_chars:
                break

        if not selected:
            # fallback: first few chunks
            joined = " ".join(self.pdf_chunks[:4])
            return joined[:max_ref_chars]

        ref_text = " ".join(selected)
        return ref_text[:max_ref_chars]

    # -------- MAIN GENERATION --------
    def generate_slide_content(self, topic, slide_title, prompt_style):
        # ===== PDF reference context (لو موجود وتم بناؤه) =====
        reference_block = ""
        if self.pdf_vectorizer is not None and self.pdf_tfidf is not None:
            ref_snippet = self._build_reference_snippet(topic, slide_title)
            if ref_snippet:
                reference_block = f"""
You have the following reference material extracted from a PDF uploaded by the lecturer.
Use it as the MAIN knowledge source and stay as close as possible to its concepts,
definitions, and terminology. Do NOT mention the PDF itself or say that you are using
reference material. Just write clean academic content.

REFERENCE MATERIAL:
\"\"\"{ref_snippet}\"\"\"


"""

        # 🔥 النسخة الذكية التي تمنع أي كلام LLM من الأساس
        prompt1 = f"""{reference_block}You are a university lecturer.

Write ONLY the final academic content for the slide titled: '{slide_title}'
in a lecture on: '{topic}'.

STRICT RULES:
- DO NOT start with phrases like: "Sure", "Certainly", "Here is", "Make sure", 
  "In this slide", "This content", "Below is", "The following".
- DO NOT explain what you are doing.
- DO NOT talk about generating content.
- DO NOT give instructions or meta-comments.
- DO NOT introduce or summarize the slide meta.
- DO NOT include helper text, transitions, or disclaimers.
- ONLY write clean academic paragraphs.

CONTENT REQUIREMENTS:
- 2–3 paragraphs.
- Each paragraph must contain 4–6 sentences.
- Use formal academic tone.
- Provide explanation + example + real application.
- No bullet points or lists.

Begin immediately with the academic content:
"""

        prompt2 = f"""{reference_block}Act only as a domain expert writing polished academic content.

Write ONLY the final paragraphs for the slide titled '{slide_title}' 
in a lecture about '{topic}'.

Forbidden:
- Any AI disclaimers.
- Any helper phrases ("Here is", "Sure", "Certainly").
- Any meta-instructions.
- Any slide descriptions.
- Any notes about what you will do.
- Any teaching/explanation about the slide structure.

Required:
- 2–3 academic paragraphs (4–6 sentences each).
- High-density, expert-level writing.
- Clear reasoning + example + real-world application.
- No lists, bullets, or headings.

Start directly with the paragraph content:
"""

        prompt = prompt1 if prompt_style == "prompt1" else prompt2

        try:
            print(f"   🤖 Generating content for: {slide_title}")
            out = self.llm(prompt, max_tokens=600)

            text = clean_text(out["choices"][0]["text"].strip())
            paragraphs = [p.strip() for p in text.split("\n") if len(p.strip()) > 40]

            return paragraphs[:3] if paragraphs else ["Content generation failed."]
        except Exception as e:
            print(f"   ❌ Error generating content: {e}")
            return [f"Error generating content: {str(e)}"]


# ==================== PPT CREATOR ====================
class PresentationCreator:
    def __init__(self, content_generator):
        self.content_generator = content_generator

    def create_presentation(self, title, num_slides, prompt_style):
        print(f"📊 Creating presentation: '{title}' ({num_slides} slides)")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = tb.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        structure = [
            "Introduction",
            "Core Explanation",
            "Important Details",
            "Examples",
            "Understanding the Concept",
            "Why It Matters",
            "Real Applications",
            "Common Mistakes",
            "Summary",
            "Final Notes"
        ][:num_slides - 1]

        for i, slide_title in enumerate(structure):
            print(f"   📝 Slide {i+1}/{len(structure)}: {slide_title}")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
            tf = tb.text_frame
            tf.text = slide_title
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True

            paragraphs = self.content_generator.generate_slide_content(
                title, slide_title, prompt_style
            )

            all_chunks = []
            for para in paragraphs:
                all_chunks.extend(split_content_if_needed(para))

            # First chunk on current slide
            cb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
            cf = cb.text_frame
            cf.word_wrap = True
            p = cf.paragraphs[0]
            p.text = all_chunks[0]
            p.font.size = Pt(18)
            p.line_spacing = 1.3

            # Extra chunks → new slides with same title
            for chunk in all_chunks[1:]:
                new_slide = prs.slides.add_slide(prs.slide_layouts[6])
                tb2 = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
                tf2 = tb2.text_frame
                tf2.text = slide_title
                tf2.paragraphs[0].font.size = Pt(32)
                tf2.paragraphs[0].font.bold = True

                cb2 = new_slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
                cf2 = cb2.text_frame
                cf2.word_wrap = True
                p2 = cf2.paragraphs[0]
                p2.text = chunk
                p2.font.size = Pt(18)
                p2.line_spacing = 1.3

        # Create temp file
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        filename = temp_file.name
        temp_file.close()
        prs.save(filename)

        # Track temp file for cleanup
        TEMP_FILES.append(filename)
        print(f"✅ Presentation saved: {filename}")

        return filename


# ==================== HTML TEMPLATE ====================
HTML_TEMPLATE = '''
ٍٍ<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>AI Presentation Creator Pro</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <link href="https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        /* ========== CSS كامل في مكان واحد ========== */
        :root {
            --primary: #667eea;
            --secondary: #764ba2;
            --accent: #ff6b6b;
            --success: #48bb78;
            --warning: #ed8936;
            --dark: #2d3748;
            --light: #f8fafc;
            --gray: #718096;
        }
        
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body { 
            font-family: 'Poppins', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif; 
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
            min-height: 100vh; 
            color: var(--dark); 
            line-height: 1.6;
            overflow-x: hidden;
        }
        
        .container { 
            max-width: 1400px; 
            margin: 0 auto; 
            padding: 20px; 
        }
        
        /* Floating Particles Animation */
        .particles {
            position: fixed;
            width: 100%;
            height: 100%;
            pointer-events: none;
            z-index: 0;
        }
        
        .particle {
            position: absolute;
            background: rgba(255, 255, 255, 0.1);
            border-radius: 50%;
            animation: float 20s infinite linear;
        }
        
        @keyframes float {
            0% { transform: translateY(0) translateX(0) rotate(0deg); opacity: 0; }
            10% { opacity: 1; }
            90% { opacity: 1; }
            100% { transform: translateY(-100vh) translateX(100vw) rotate(720deg); opacity: 0; }
        }
        
        /* Header Styles */
        .header { 
            text-align: center; 
            color: white; 
            margin-bottom: 40px; 
            padding: 40px 0; 
            position: relative;
            z-index: 1;
        }
        
        .logo { 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            gap: 20px; 
            margin-bottom: 20px;
            animation: fadeInDown 1s ease;
        }
        
        .logo-icon {
            font-size: 3.5rem;
            color: #ffd700;
            animation: pulse 2s infinite;
        }
        
        .logo h1 { 
            font-size: 3.2rem; 
            font-weight: 700; 
            text-shadow: 0 4px 20px rgba(0,0,0,0.3);
            background: linear-gradient(45deg, #fff, #ffd700);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }
        
        .tagline { 
            font-size: 1.3rem; 
            opacity: 0.9; 
            max-width: 700px; 
            margin: 0 auto;
            animation: fadeInUp 1s ease 0.3s both;
        }
        
        /* Main Content */
        .main-content {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 30px;
            margin-bottom: 40px;
            position: relative;
            z-index: 1;
        }
        
        @media (max-width: 1024px) {
            .main-content {
                grid-template-columns: 1fr;
            }
        }
        
        /* Cards */
        .card { 
            background: rgba(255, 255, 255, 0.95); 
            border-radius: 25px; 
            padding: 40px; 
            box-shadow: 0 25px 70px rgba(0,0,0,0.3); 
            margin-bottom: 30px; 
            transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
            backdrop-filter: blur(10px);
            border: 1px solid rgba(255, 255, 255, 0.2);
            animation: slideUp 0.8s ease;
        }
        
        .card:hover { 
            transform: translateY(-10px) scale(1.02);
            box-shadow: 0 35px 80px rgba(0,0,0,0.4);
        }
        
        .card-header { 
            margin-bottom: 35px; 
            border-bottom: 2px solid var(--light); 
            padding-bottom: 25px; 
        }
        
        .card-header h2 { 
            color: var(--dark); 
            font-size: 2rem; 
            display: flex; 
            align-items: center; 
            gap: 15px;
            animation: fadeInLeft 0.8s ease;
        }
        
        .card-header h2 i {
            color: var(--primary);
            animation: spinSlow 10s linear infinite;
        }
        
        .card-header p { 
            color: var(--gray); 
            margin-top: 10px;
            font-size: 1.1rem;
        }
        
        /* Form Styles */
        .form-group { 
            margin-bottom: 30px; 
            animation: fadeInRight 0.8s ease;
        }
        
        .form-group label { 
            display: block; 
            color: var(--dark); 
            font-weight: 600; 
            margin-bottom: 12px; 
            font-size: 1.15rem; 
            display: flex; 
            align-items: center; 
            gap: 12px;
        }
        
        input[type="text"] { 
            width: 100%; 
            padding: 18px 24px; 
            border: 2px solid #e2e8f0; 
            border-radius: 15px; 
            font-size: 1.1rem; 
            transition: all 0.3s ease; 
            background: var(--light);
            font-family: 'Poppins', sans-serif;
        }
        
        input[type="text"]:focus { 
            outline: none; 
            border-color: var(--primary); 
            background: white; 
            box-shadow: 0 0 0 4px rgba(102, 126, 234, 0.2);
            transform: translateY(-2px);
        }
        
        .hint { 
            color: var(--gray); 
            font-size: 0.95rem; 
            margin-top: 10px; 
            display: flex; 
            align-items: center; 
            gap: 8px;
            opacity: 0.8;
        }
        
        /* File Upload Area */
        .file-upload-area {
            border: 3px dashed var(--primary);
            border-radius: 20px;
            padding: 40px;
            text-align: center;
            background: rgba(102, 126, 234, 0.05);
            cursor: pointer;
            transition: all 0.3s ease;
            margin-bottom: 20px;
            position: relative;
            overflow: hidden;
        }
        
        .file-upload-area:hover {
            background: rgba(102, 126, 234, 0.1);
            border-color: var(--secondary);
            transform: translateY(-5px);
        }
        
        .file-upload-area i {
            font-size: 4rem;
            color: var(--primary);
            margin-bottom: 20px;
            display: block;
            animation: bounce 2s infinite;
        }
        
        .file-upload-area h3 {
            color: var(--dark);
            margin-bottom: 10px;
            font-size: 1.5rem;
        }
        
        .file-upload-area p {
            color: var(--gray);
            margin-bottom: 20px;
        }
        
        .file-input {
            display: none;
        }
        
        .file-info {
            display: none;
            background: var(--success);
            color: white;
            padding: 15px;
            border-radius: 10px;
            margin-top: 15px;
            animation: fadeIn 0.5s ease;
        }
        
        .file-info.show {
            display: block;
        }
        
        .file-info i {
            font-size: 1.2rem;
            margin-right: 10px;
        }
        
        /* Range Slider */
        .range-container { 
            background: var(--light); 
            padding: 30px; 
            border-radius: 15px; 
            border: 2px solid #e2e8f0;
            position: relative;
            overflow: hidden;
        }
        
        .range-container:before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: linear-gradient(90deg, var(--primary), var(--secondary));
            opacity: 0.1;
            z-index: 0;
        }
        
        input[type="range"] { 
            width: 100%; 
            height: 12px; 
            -webkit-appearance: none; 
            appearance: none; 
            background: linear-gradient(90deg, var(--primary) 0%, var(--secondary) 100%); 
            border-radius: 6px; 
            outline: none;
            position: relative;
            z-index: 1;
        }
        
        input[type="range"]::-webkit-slider-thumb { 
            -webkit-appearance: none; 
            width: 28px; 
            height: 28px; 
            background: white; 
            border-radius: 50%; 
            border: 4px solid var(--primary); 
            cursor: pointer; 
            box-shadow: 0 4px 15px rgba(0,0,0,0.3);
            transition: all 0.3s ease;
        }
        
        input[type="range"]::-webkit-slider-thumb:hover {
            transform: scale(1.2);
            box-shadow: 0 6px 20px rgba(0,0,0,0.4);
        }
        
        .range-value { 
            text-align: center; 
            margin-top: 20px; 
            font-size: 1.4rem; 
            font-weight: 600; 
            color: var(--dark);
            position: relative;
            z-index: 1;
        }
        
        .range-value span { 
            color: var(--primary); 
            font-size: 1.8rem;
            font-weight: 700;
            text-shadow: 0 2px 10px rgba(102, 126, 234, 0.3);
        }
        
        /* Style Options */
        .style-options { 
            display: grid; 
            grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); 
            gap: 20px; 
            margin-top: 15px;
        }
        
        .style-option input { 
            display: none; 
        }
        
        .style-option label { 
            display: block; 
            padding: 25px; 
            border: 2px solid #e2e8f0; 
            border-radius: 15px; 
            cursor: pointer; 
            transition: all 0.3s ease; 
            background: var(--light);
            position: relative;
            overflow: hidden;
        }
        
        .style-option label:before {
            content: '';
            position: absolute;
            top: 0;
            left: -100%;
            width: 100%;
            height: 100%;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.4), transparent);
            transition: 0.5s;
        }
        
        .style-option label:hover:before {
            left: 100%;
        }
        
        .style-option input:checked + label { 
            border-color: var(--primary); 
            background: linear-gradient(135deg, rgba(102, 126, 234, 0.15) 0%, rgba(118, 75, 162, 0.15) 100%); 
            transform: translateY(-5px) scale(1.03);
            box-shadow: 0 10px 30px rgba(102, 126, 234, 0.2);
        }
        
        .style-option label i { 
            font-size: 1.8rem; 
            color: var(--primary); 
            margin-bottom: 15px;
            transition: all 0.3s ease;
        }
        
        .style-option input:checked + label i {
            transform: rotate(360deg);
            color: var(--secondary);
        }
        
        .style-option label span { 
            display: block; 
            font-weight: 600; 
            color: var(--dark); 
            margin-bottom: 8px;
            font-size: 1.2rem;
        }
        
        .style-option label small { 
            color: var(--gray); 
            font-size: 0.95rem; 
            display: block; 
            line-height: 1.5;
        }
        
        /* Generate Button */
        .btn-generate { 
            width: 100%; 
            padding: 25px; 
            background: linear-gradient(135deg, var(--primary) 0%, var(--secondary) 100%); 
            color: white; 
            border: none; 
            border-radius: 15px; 
            font-size: 1.3rem; 
            font-weight: 600; 
            cursor: pointer; 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            gap: 15px; 
            transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
            position: relative;
            overflow: hidden;
            font-family: 'Poppins', sans-serif;
            text-transform: uppercase;
            letter-spacing: 1px;
            margin-top: 10px;
        }
        
        .btn-generate:before {
            content: '';
            position: absolute;
            top: 0;
            left: -100%;
            width: 100%;
            height: 100%;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3), transparent);
            transition: 0.5s;
        }
        
        .btn-generate:hover:before {
            left: 100%;
        }
        
        .btn-generate:hover { 
            transform: translateY(-5px) scale(1.02); 
            box-shadow: 0 20px 40px rgba(102, 126, 234, 0.5);
            letter-spacing: 2px;
        }
        
        .btn-generate:active { 
            transform: translateY(0) scale(0.98); 
        }
        
        .btn-generate i {
            font-size: 1.5rem;
            transition: all 0.3s ease;
        }
        
        .btn-generate:hover i {
            transform: rotate(45deg) scale(1.2);
        }
        
        /* Preview Card */
        .preview-card { 
            background: linear-gradient(135deg, var(--dark) 0%, #4a5568 100%); 
            color: white;
            animation: slideUp 0.8s ease 0.2s both;
        }
        
        .preview-card .card-header h2 { 
            color: white; 
        }
        
        .preview-content { 
            display: grid; 
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr)); 
            gap: 30px; 
        }
        
        .preview-slide { 
            background: rgba(255, 255, 255, 0.15); 
            border-radius: 15px; 
            padding: 30px; 
            backdrop-filter: blur(10px); 
            border: 1px solid rgba(255, 255, 255, 0.2);
            transition: all 0.4s ease;
            position: relative;
            overflow: hidden;
        }
        
        .preview-slide:before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 5px;
            background: linear-gradient(90deg, var(--primary), var(--accent));
        }
        
        .preview-slide:hover { 
            transform: translateY(-10px) scale(1.05);
            box-shadow: 0 20px 40px rgba(0,0,0,0.3);
        }
        
        .slide-header { 
            font-weight: 600; 
            color: #ffd700; 
            margin-bottom: 20px; 
            font-size: 1.2rem;
            display: flex;
            align-items: center;
            gap: 10px;
        }
        
        .slide-content h3, .slide-content h4 { 
            margin-bottom: 20px; 
            color: white; 
        }
        
        .slide-content ul { 
            list-style: none; 
            padding-left: 0; 
        }
        
        .slide-content li { 
            margin-bottom: 12px; 
            padding-left: 30px; 
            position: relative; 
            transition: all 0.3s ease;
        }
        
        .slide-content li:hover {
            transform: translateX(10px);
            color: #ffd700;
        }
        
        .slide-content li:before { 
            content: "✨"; 
            position: absolute; 
            left: 0; 
            color: var(--accent); 
            font-weight: bold;
            animation: twinkle 2s infinite;
        }
        
        /* Footer */
        .footer { 
            text-align: center; 
            color: white; 
            padding: 40px 0; 
            margin-top: 60px; 
            border-top: 1px solid rgba(255, 255, 255, 0.2); 
            position: relative;
            z-index: 1;
        }
        
        .footer-content p { 
            margin-bottom: 20px; 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            gap: 15px; 
            font-size: 1.1rem;
            animation: fadeIn 1s ease 0.5s both;
        }
        
        .tech-stack { 
            display: flex; 
            justify-content: center; 
            gap: 20px; 
            flex-wrap: wrap; 
        }
        
        .tech-tag { 
            background: rgba(255, 255, 255, 0.15); 
            padding: 12px 24px; 
            border-radius: 25px; 
            font-size: 1rem; 
            backdrop-filter: blur(5px);
            transition: all 0.3s ease;
            border: 1px solid rgba(255, 255, 255, 0.1);
            animation: fadeInUp 0.5s ease;
        }
        
        .tech-tag:hover {
            background: rgba(255, 255, 255, 0.25);
            transform: translateY(-3px);
        }
        
        /* Modal */
        .modal { 
            display: none; 
            position: fixed; 
            top: 0; 
            left: 0; 
            width: 100%; 
            height: 100%; 
            background: rgba(0, 0, 0, 0.9); 
            z-index: 10000; 
            backdrop-filter: blur(10px); 
            align-items: center; 
            justify-content: center;
            animation: fadeIn 0.3s ease;
        }
        
        .modal-content { 
            background: linear-gradient(135deg, #fff 0%, #f8fafc 100%); 
            padding: 60px; 
            border-radius: 25px; 
            text-align: center; 
            max-width: 600px; 
            width: 90%; 
            box-shadow: 0 40px 80px rgba(0,0,0,0.5);
            animation: modalAppear 0.5s cubic-bezier(0.4, 0, 0.2, 1);
            position: relative;
            overflow: hidden;
        }
        
        .modal-content:before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 5px;
            background: linear-gradient(90deg, var(--primary), var(--secondary), var(--accent));
        }
        
        @keyframes modalAppear { 
            from { 
                opacity: 0; 
                transform: translateY(-50px) scale(0.9) rotateX(10deg); 
            } 
            to { 
                opacity: 1; 
                transform: translateY(0) scale(1) rotateX(0); 
            } 
        }
        
        .loading-spinner .spinner { 
            width: 100px; 
            height: 100px; 
            border: 8px solid #f3f3f3; 
            border-top: 8px solid var(--primary); 
            border-radius: 50%; 
            animation: spin 1s linear infinite; 
            margin: 0 auto 40px;
            position: relative;
        }
        
        .loading-spinner .spinner:after {
            content: '';
            position: absolute;
            top: -8px;
            left: -8px;
            right: -8px;
            bottom: -8px;
            border: 8px solid transparent;
            border-radius: 50%;
            border-top: 8px solid var(--secondary);
            animation: spin 2s linear infinite reverse;
        }
        
        .loading-spinner .spinner:before {
            content: '';
            position: absolute;
            top: -16px;
            left: -16px;
            right: -16px;
            bottom: -16px;
            border: 8px solid transparent;
            border-radius: 50%;
            border-top: 8px solid var(--accent);
            animation: spin 3s linear infinite;
        }
        
        @keyframes spin { 
            0% { 
                transform: rotate(0deg); 
            } 
            100% { 
                transform: rotate(360deg); 
            } 
        }
        
        .loading-spinner h3 { 
            margin-bottom: 20px; 
            color: var(--dark); 
            font-size: 1.8rem;
        }
        
        .loading-spinner p { 
            color: var(--gray); 
            margin-bottom: 30px;
            font-size: 1.1rem;
        }
        
        .progress-container { 
            height: 12px; 
            background: #f0f0f0; 
            border-radius: 6px; 
            margin: 40px 0; 
            overflow: hidden; 
            position: relative;
        }
        
        .progress-bar { 
            height: 100%; 
            background: linear-gradient(90deg, var(--primary) 0%, var(--secondary) 100%); 
            width: 0%; 
            transition: width 0.5s cubic-bezier(0.4, 0, 0.2, 1);
            border-radius: 6px;
            position: relative;
            overflow: hidden;
        }
        
        .progress-bar:after {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.4), transparent);
            animation: shimmer 2s infinite;
        }
        
        .loading-details { 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            gap: 15px; 
            color: var(--primary); 
            font-weight: 600;
            font-size: 1.1rem;
            animation: pulse 2s infinite;
        }
        
        /* Notifications */
        .notification { 
            position: fixed; 
            top: 30px; 
            right: 30px; 
            padding: 20px 25px; 
            border-radius: 12px; 
            display: flex; 
            align-items: center; 
            gap: 15px; 
            box-shadow: 0 10px 30px rgba(0,0,0,0.2); 
            z-index: 10000; 
            animation: slideInRight 0.5s cubic-bezier(0.4, 0, 0.2, 1);
            max-width: 450px;
            backdrop-filter: blur(10px);
            border: 1px solid rgba(255,255,255,0.2);
        }
        
        .notification.error { 
            background: linear-gradient(135deg, #fed7d7 0%, #feb2b2 100%); 
            color: #9b2c2c; 
        }
        
        .notification.success { 
            background: linear-gradient(135deg, #c6f6d5 0%, #9ae6b4 100%); 
            color: #22543d; 
        }
        
        .notification-close { 
            background: none; 
            border: none; 
            cursor: pointer; 
            padding: 0; 
            margin-left: 15px; 
            color: inherit;
            font-size: 1.2rem;
            transition: all 0.3s ease;
        }
        
        .notification-close:hover {
            transform: rotate(90deg) scale(1.2);
        }
        
        /* Animations */
        @keyframes fadeIn {
            from { opacity: 0; }
            to { opacity: 1; }
        }
        
        @keyframes fadeInDown {
            from { 
                opacity: 0; 
                transform: translateY(-30px); 
            }
            to { 
                opacity: 1; 
                transform: translateY(0); 
            }
        }
        
        @keyframes fadeInUp {
            from { 
                opacity: 0; 
                transform: translateY(30px); 
            }
            to { 
                opacity: 1; 
                transform: translateY(0); 
            }
        }
        
        @keyframes fadeInLeft {
            from { 
                opacity: 0; 
                transform: translateX(-30px); 
            }
            to { 
                opacity: 1; 
                transform: translateX(0); 
            }
        }
        
        @keyframes fadeInRight {
            from { 
                opacity: 0; 
                transform: translateX(30px); 
            }
            to { 
                opacity: 1; 
                transform: translateX(0); 
            }
        }
        
        @keyframes slideUp {
            from { 
                opacity: 0; 
                transform: translateY(50px) scale(0.95); 
            }
            to { 
                opacity: 1; 
                transform: translateY(0) scale(1); 
            }
        }
        
        @keyframes slideInRight {
            from { 
                transform: translateX(100%); 
                opacity: 0; 
            } 
            to { 
                transform: translateX(0); 
                opacity: 1; 
            } 
        }
        
        @keyframes slideOutRight { 
            from { 
                transform: translateX(0); 
                opacity: 1; 
            } 
            to { 
                transform: translateX(100%); 
                opacity: 0; 
            } 
        }
        
        @keyframes pulse {
            0%, 100% { transform: scale(1); }
            50% { transform: scale(1.05); }
        }
        
        @keyframes bounce {
            0%, 100% { transform: translateY(0); }
            50% { transform: translateY(-10px); }
        }
        
        @keyframes twinkle {
            0%, 100% { opacity: 1; }
            50% { opacity: 0.5; }
        }
        
        @keyframes spinSlow {
            from { transform: rotate(0deg); }
            to { transform: rotate(360deg); }
        }
        
        @keyframes shimmer {
            0% { transform: translateX(-100%); }
            100% { transform: translateX(100%); }
        }
        
        /* Character Counter */
        .char-counter { 
            text-align: right; 
            margin-top: 8px; 
            font-size: 0.95rem; 
            font-weight: 500;
            transition: all 0.3s ease;
        }
        
        /* Responsive Design */
        @media (max-width: 768px) {
            .container { 
                padding: 15px; 
            }
            
            .logo { 
                flex-direction: column; 
                text-align: center; 
                gap: 15px; 
            }
            
            .logo h1 { 
                font-size: 2.2rem; 
            }
            
            .logo-icon {
                font-size: 2.8rem;
            }
            
            .tagline { 
                font-size: 1.1rem; 
            }
            
            .card { 
                padding: 25px; 
            }
            
            .style-options { 
                grid-template-columns: 1fr; 
            }
            
            .modal-content { 
                padding: 30px 20px; 
            }
            
            .preview-content {
                grid-template-columns: 1fr;
            }
        }
        
        @media (max-width: 480px) {
            .logo h1 { 
                font-size: 1.8rem; 
            }
            
            .card-header h2 { 
                font-size: 1.5rem; 
            }
            
            .btn-generate { 
                padding: 20px; 
                font-size: 1.1rem; 
            }
            
            .file-upload-area {
                padding: 25px;
            }
            
            .tech-tag {
                padding: 10px 20px;
                font-size: 0.9rem;
            }
        }
        
        /* Particle Effects */
        .confetti {
            position: fixed;
            width: 10px;
            height: 10px;
            background-color: var(--primary);
            border-radius: 50%;
            animation: confetti-fall 3s linear forwards;
        }
        
        @keyframes confetti-fall {
            0% {
                transform: translateY(-100vh) rotate(0deg);
                opacity: 1;
            }
            100% {
                transform: translateY(100vh) rotate(720deg);
                opacity: 0;
            }
        }
        
        /* Typewriter Effect */
        .typewriter {
            overflow: hidden;
            border-right: .15em solid var(--primary);
            white-space: nowrap;
            margin: 0 auto;
            letter-spacing: .15em;
            animation: typing 3.5s steps(40, end), blink-caret .75s step-end infinite;
        }
        
        @keyframes typing {
            from { width: 0 }
            to { width: 100% }
        }
        
        @keyframes blink-caret {
            from, to { border-color: transparent }
            50% { border-color: var(--primary) }
        }
    </style>
</head>
<body>
    <!-- Floating Particles -->
    <div class="particles" id="particles"></div>
    
    <div class="container">
        <header class="header">
            <div class="logo">
                <i class="fas fa-chalkboard-teacher logo-icon"></i>
                <h1 class="typewriter">AI Presentation Creator Pro</h1>
            </div>
            <p class="tagline">Generate professional academic presentations with AI intelligence. Now enhanced with PDF learning!</p>
        </header>
        
        <main class="main-content">
            <!-- Left Column: Form -->
            <div class="card">
                <div class="card-header">
                    <h2><i class="fas fa-magic"></i> Create Your Presentation</h2>
                    <p>Fill in the details below to generate an intelligent PowerPoint presentation</p>
                </div>
                <form id="presentationForm" class="form" enctype="multipart/form-data">
                    <div class="form-group">
                        <label for="title"><i class="fas fa-heading"></i> Presentation Title</label>
                        <input type="text" id="title" name="title" placeholder="Enter your presentation topic (e.g., Machine Learning Fundamentals)" required minlength="3">
                        <div class="hint">
                            <i class="fas fa-lightbulb"></i>
                            Minimum 3 characters. Be specific for better results!
                        </div>
                        <div class="char-counter" id="charCounter">0 characters</div>
                    </div>
                    
                    <div class="form-group">
                        <label for="pdfFile"><i class="fas fa-file-pdf"></i> Upload Reference PDF (Optional)</label>
                        <div class="file-upload-area" id="fileUploadArea">
                            <i class="fas fa-cloud-upload-alt"></i>
                            <h3>Drag & Drop PDF Here</h3>
                            <p>Or click to browse files</p>
                            <p class="hint">Upload a PDF to enhance AI accuracy with specific content</p>
                            <input type="file" id="pdfFile" name="pdfFile" class="file-input" accept=".pdf">
                            <button type="button" class="btn-generate" style="padding: 15px; max-width: 700px; margin: 0 auto;">
                                <i class="fas fa-folder-open"></i>
                                <span>Browse PDF</span>
                            </button>
                        </div>
                        <div class="file-info" id="fileInfo"></div>
                    </div>
                    
                    <div class="form-group">
                        <label for="slides"><i class="fas fa-sliders-h"></i> Number of Slides</label>
                        <div class="range-container">
                            <input type="range" id="slides" name="slides" min="3" max="12" value="8">
                            <div class="range-value"><span id="slideValue">8</span> slides</div>
                        </div>
                        <div class="hint">
                            <i class="fas fa-chart-line"></i>
                            Recommended: 8-10 slides for optimal content density
                        </div>
                    </div>
                    
                    <div class="form-group">
                        <label for="promptStyle"><i class="fas fa-language"></i> Content Style</label>
                        <div class="style-options">
                            <div class="style-option">
                                <input type="radio" id="prompt1" name="prompt_style" value="prompt1" checked>
                                <label for="prompt1">
                                    <i class="fas fa-graduation-cap"></i>
                                    <span>Standard Academic</span>
                                    <small>Clear, structured content for general lectures and presentations</small>
                                </label>
                            </div>
                            <div class="style-option">
                                <input type="radio" id="prompt2" name="prompt_style" value="prompt2">
                                <label for="prompt2">
                                    <i class="fas fa-flask"></i>
                                    <span>Advanced Expert</span>
                                    <small>Dense, technical content for specialized topics and research</small>
                                </label>
                            </div>
                        </div>
                    </div>
                    
                    <div class="form-group">
                        <button type="submit" id="generateBtn" class="btn-generate">
                            <i class="fas fa-rocket"></i>
                            <span>Generate Intelligent Presentation</span>
                        </button>
                        <div class="disclaimer">
                            <i class="fas fa-info-circle"></i>
                            First generation may take 1-2 minutes as AI model loads. PDF upload enhances accuracy!
                        </div>
                    </div>
                </form>
            </div>
            
            <!-- Right Column: Preview -->
            <div class="card preview-card">
                <div class="card-header">
                    <h2><i class="fas fa-eye"></i> What You'll Get</h2>
                </div>
                <div class="preview-content">
                    <div class="preview-slide">
                        <div class="slide-header">
                            <i class="fas fa-star"></i>
                            Title Slide
                        </div>
                        <div class="slide-content">
                            <h3>YOUR PRESENTATION TITLE</h3>
                            <p>Professional title slide with elegant typography and design</p>
                        </div>
                    </div>
                    <div class="preview-slide">
                        <div class="slide-header">
                            <i class="fas fa-file-alt"></i>
                            Content Slides
                        </div>
                        <div class="slide-content">
                            <h4>AI-Enhanced Content</h4>
                            <ul>
                                <li>2-3 well-structured academic paragraphs</li>
                                <li>PDF-enhanced accuracy and specificity</li>
                                <li>Real-world applications & examples</li>
                                <li>Professional academic formatting</li>
                                <li>Citations from reference material</li>
                            </ul>
                        </div>
                    </div>
                    <div class="preview-slide">
                        <div class="slide-header">
                            <i class="fas fa-brain"></i>
                            AI Intelligence
                        </div>
                        <div class="slide-content">
                            <h4>Enhanced with PDF Learning</h4>
                            <ul>
                                <li>Context-aware content generation</li>
                                <li>PDF reference integration</li>
                                <li>Domain-specific terminology</li>
                                <li>Accurate technical details</li>
                                <li>Citation of source material</li>
                            </ul>
                        </div>
                    </div>
                </div>
            </div>
        </main>
        
        <footer class="footer">
            <div class="footer-content">
                <p><i class="fas fa-brain"></i> Powered by Advanced AI & Local LLM Technology</p>
                <p class="tech-stack">
                    <span class="tech-tag">Flask</span>
                    <span class="tech-tag">LLaMA.cpp</span>
                    <span class="tech-tag">python-pptx</span>
                    <span class="tech-tag">PyPDF2</span>
                    <span class="tech-tag">PDF Learning</span>
                    <span class="tech-tag">Advanced AI</span>
                </p>
                <p class="hint" style="color: rgba(255,255,255,0.7); margin-top: 20px;">
                    <i class="fas fa-lightbulb"></i>
                    Upload PDFs to train the AI on specific content for enhanced accuracy!
                </p>
            </div>
        </footer>
    </div>
    
    <!-- Loading Modal -->
    <div id="loadingModal" class="modal">
        <div class="modal-content">
            <div class="loading-spinner">
                <div class="spinner"></div>
                <h3>Generating Your Intelligent Presentation</h3>
                <p>AI is processing your request. This may take a moment. Please don't close this window.</p>
                <div class="progress-container">
                    <div class="progress-bar" id="progressBar"></div>
                </div>
                <p class="loading-details" id="loadingDetails">
                    <i class="fas fa-cog fa-spin"></i>
                    Initializing AI model and processing PDF...
                </p>
            </div>
        </div>
    </div>

    <!-- JavaScript كامل في مكان واحد -->
    <script>
        // ========== JavaScript كامل ==========
        document.addEventListener('DOMContentLoaded', function() {
            console.log('🎉 Presentation Creator Pro loaded successfully!');
            
            // عناصر DOM الرئيسية
            const form = document.getElementById('presentationForm');
            const slidesInput = document.getElementById('slides');
            const slideValue = document.getElementById('slideValue');
            const loadingModal = document.getElementById('loadingModal');
            const progressBar = document.getElementById('progressBar');
            const loadingDetails = document.getElementById('loadingDetails');
            const titleInput = document.getElementById('title');
            const charCounter = document.getElementById('charCounter');
            const fileUploadArea = document.getElementById('fileUploadArea');
            const pdfFileInput = document.getElementById('pdfFile');
            const fileInfo = document.getElementById('fileInfo');
            const browseBtn = fileUploadArea.querySelector('button');
            const generateBtn = document.getElementById('generateBtn');
            
            // حالة التطبيق
            let uploadedPDF = null;
            let isLoading = false;
            let particlesCreated = false;
            
            // ========== 1. إنشاء الجسيمات المتحركة ==========
            function createParticles() {
                if (particlesCreated) return;
                
                const particlesContainer = document.getElementById('particles');
                const particleCount = 50;
                
                for (let i = 0; i < particleCount; i++) {
                    const particle = document.createElement('div');
                    particle.className = 'particle';
                    
                    // خصائص عشوائية
                    const size = Math.random() * 10 + 5;
                    const duration = Math.random() * 10 + 20;
                    const delay = Math.random() * 20;
                    const color = Math.random() > 0.5 ? 'rgba(255, 255, 255, 0.1)' : 'rgba(255, 215, 0, 0.1)';
                    
                    particle.style.cssText = `
                        width: ${size}px;
                        height: ${size}px;
                        left: ${Math.random() * 100}%;
                        top: ${Math.random() * 100}%;
                        background: ${color};
                        animation-delay: ${delay}s;
                        animation-duration: ${duration}s;
                    `;
                    
                    particlesContainer.appendChild(particle);
                }
                
                particlesCreated = true;
                console.log(`✨ Created ${particleCount} floating particles`);
            }
            
            // ========== 2. تفعيل الجسيمات ==========
            createParticles();
            
            // ========== 3. تحديث شريط الشرائح ==========
            slidesInput.addEventListener('input', () => {
                slideValue.textContent = slidesInput.value;
                slideValue.style.transform = 'scale(1.2)';
                slideValue.style.color = '#667eea';
                
                setTimeout(() => {
                    slideValue.style.transform = 'scale(1)';
                    slideValue.style.color = '';
                }, 300);
            });
            
            // ========== 4. عداد الحروف ==========
            function updateCharCounter() {
                const length = titleInput.value.length;
                charCounter.textContent = `${length} characters`;
                
                if (length < 3) {
                    charCounter.style.color = '#f56565';
                    charCounter.style.fontWeight = 'bold';
                } else if (length < 10) {
                    charCounter.style.color = '#ed8936';
                    charCounter.style.fontWeight = '600';
                } else {
                    charCounter.style.color = '#48bb78';
                    charCounter.style.fontWeight = '600';
                }
            }
            
            titleInput.addEventListener('input', updateCharCounter);
            updateCharCounter();
            
            // ========== 5. معالجة رفع الملفات ==========
            browseBtn.addEventListener('click', () => {
                pdfFileInput.click();
            });
            
            fileUploadArea.addEventListener('click', (e) => {
                if (e.target !== browseBtn && e.target !== pdfFileInput) {
                    pdfFileInput.click();
                }
            });
            
            // Drag and Drop
            fileUploadArea.addEventListener('dragover', (e) => {
                e.preventDefault();
                fileUploadArea.style.transform = 'scale(1.05)';
                fileUploadArea.style.borderColor = '#48bb78';
                fileUploadArea.style.boxShadow = '0 10px 30px rgba(72, 187, 120, 0.3)';
            });
            
            fileUploadArea.addEventListener('dragleave', () => {
                fileUploadArea.style.transform = 'scale(1)';
                fileUploadArea.style.borderColor = 'var(--primary)';
                fileUploadArea.style.boxShadow = 'none';
            });
            
            fileUploadArea.addEventListener('drop', (e) => {
                e.preventDefault();
                fileUploadArea.style.transform = 'scale(1)';
                fileUploadArea.style.borderColor = 'var(--primary)';
                fileUploadArea.style.boxShadow = 'none';
                
                if (e.dataTransfer.files.length) {
                    pdfFileInput.files = e.dataTransfer.files;
                    handleFileSelection(e.dataTransfer.files[0]);
                }
            });
            
            pdfFileInput.addEventListener('change', (e) => {
                if (e.target.files.length) {
                    handleFileSelection(e.target.files[0]);
                }
            });
            
            function handleFileSelection(file) {
                if (file.type !== 'application/pdf') {
                    showError('Please select a PDF file');
                    return;
                }
                
                if (file.size > 10 * 1024 * 1024) {
                    showError('File size must be less than 10MB');
                    return;
                }
                
                uploadedPDF = file;
                fileInfo.innerHTML = `
                    <i class="fas fa-check-circle"></i>
                    <strong>${file.name}</strong> (${(file.size / 1024 / 1024).toFixed(2)} MB)
                    <button type="button" id="removeFile" style="background: none; border: none; color: white; margin-left: 10px; cursor: pointer;">
                        <i class="fas fa-times"></i>
                    </button>
                `;
                fileInfo.classList.add('show');
                
                // إضافة حدث إزالة الملف
                document.getElementById('removeFile').addEventListener('click', (e) => {
                    e.stopPropagation();
                    uploadedPDF = null;
                    pdfFileInput.value = '';
                    fileInfo.classList.remove('show');
                    showSuccess('PDF file removed');
                });
                
                showSuccess(`PDF uploaded: ${file.name}`);
            }
            
            // ========== 6. تأثيرات الـ Preview Slides ==========
            document.querySelectorAll('.preview-slide').forEach((slide, index) => {
                slide.addEventListener('mouseenter', function() {
                    this.style.transform = 'translateY(-15px) scale(1.08)';
                    this.style.transition = 'all 0.4s cubic-bezier(0.4, 0, 0.2, 1)';
                    this.style.boxShadow = '0 20px 40px rgba(0,0,0,0.4)';
                });
                
                slide.addEventListener('mouseleave', function() {
                    this.style.transform = 'translateY(0) scale(1)';
                    this.style.boxShadow = '0 10px 30px rgba(0,0,0,0.2)';
                });
                
                // تأثير عند التحميل
                setTimeout(() => {
                    slide.style.animation = 'slideUp 0.8s ease';
                }, index * 100);
            });
            
            // ========== 7. معالجة إرسال النموذج ==========
            form.addEventListener('submit', async function(e) {
                e.preventDefault();
                
                if (isLoading) {
                    showError('Please wait for the current generation to complete');
                    return;
                }
                
                // التحقق من المدخلات
                const title = titleInput.value.trim();
                if (title.length < 3) {
                    showError('Please enter a title with at least 3 characters');
                    titleInput.focus();
                    addShakeEffect(titleInput);
                    return;
                }
                
                // بدء التحميل
                startLoading();
                isLoading = true;
                
                try {
                    // إعداد البيانات
                    const formData = new FormData();
                    formData.append('title', title);
                    formData.append('slides', slidesInput.value);
                    formData.append('prompt_style', document.querySelector('input[name="prompt_style"]:checked').value);
                    
                    if (uploadedPDF) {
                        formData.append('pdf_file', uploadedPDF);
                    }
                    
                    // محاكاة التقدم
                    simulateProgress(0, 30, 100, 'Starting AI engine...');
                    
                    // إرسال الطلب
                    updateLoadingMessage('🚀 Connecting to AI brain...', 'rocket');
                    
                    const response = await fetch('/generate', {
                        method: 'POST',
                        body: formData
                    });
                    
                    // تحقق من الاستجابة
                    if (!response.ok) {
                        const errorData = await response.json().catch(() => ({}));
                        throw new Error(errorData.error || `Server error: ${response.status}`);
                    }
                    
                    // تحديث حالة التحميل
                    updateLoadingMessage('📄 Processing PDF content...', 'file-pdf');
                    updateProgress(40);
                    
                    // انتظر قليلاً لمحاكاة المعالجة
                    await new Promise(resolve => setTimeout(resolve, 1500));
                    
                    updateLoadingMessage('🤖 AI is generating smart content...', 'robot');
                    updateProgress(60);
                    
                    await new Promise(resolve => setTimeout(resolve, 2000));
                    
                    updateLoadingMessage('🎨 Designing beautiful slides...', 'palette');
                    updateProgress(80);
                    
                    // تحميل الملف
                    const blob = await response.blob();
                    const url = window.URL.createObjectURL(blob);
                    const a = document.createElement('a');
                    a.href = url;
                    a.download = `${title.replace(/[^a-z0-9]/gi, '_')}_presentation.pptx`;
                    document.body.appendChild(a);
                    a.click();
                    window.URL.revokeObjectURL(url);
                    document.body.removeChild(a);
                    
                    // النجاح
                    updateLoadingMessage('✅ Presentation ready! Downloading...', 'check-circle');
                    updateProgress(100);
                    
                    // إظهار تأثير النجاح
                    showConfetti();
                    
                    setTimeout(() => {
                        stopLoading();
                        showSuccess('🎉 Presentation generated and downloaded successfully!');
                        
                        // إعادة تعيين النموذج
                        resetForm();
                        
                        // إعادة إنشاء الجسيمات
                        createParticles();
                        
                    }, 2000);
                    
                } catch (error) {
                    console.error('Generation error:', error);
                    
                    updateLoadingMessage(`❌ Error: ${error.message}`, 'exclamation-triangle');
                    updateProgress(0);
                    progressBar.style.background = 'linear-gradient(90deg, #f56565 0%, #c53030 100%)';
                    
                    setTimeout(() => {
                        stopLoading();
                        showError(`Failed to generate presentation: ${error.message}`);
                    }, 2000);
                    
                } finally {
                    isLoading = false;
                }
            });
            
            // ========== 8. دوال الـ Loading ==========
            function startLoading() {
                loadingModal.style.display = 'flex';
                progressBar.style.width = '0%';
                progressBar.style.background = 'linear-gradient(90deg, var(--primary) 0%, var(--secondary) 100%)';
                generateBtn.disabled = true;
                generateBtn.innerHTML = '<i class="fas fa-spinner fa-spin"></i> Processing...';
            }
            
            function stopLoading() {
                loadingModal.style.display = 'none';
                generateBtn.disabled = false;
                generateBtn.innerHTML = '<i class="fas fa-rocket"></i> Generate Intelligent Presentation';
            }
            
            function updateProgress(percent) {
                progressBar.style.width = `${percent}%`;
            }
            
            function updateLoadingMessage(message, icon = 'cog') {
                loadingDetails.innerHTML = `<i class="fas fa-${icon} ${icon !== 'check-circle' ? 'fa-spin' : ''}"></i> ${message}`;
            }
            
            function simulateProgress(start, end, speed, message = 'Processing...') {
                let progress = start;
                const interval = setInterval(() => {
                    progress += 1;
                    if (progress > end) {
                        clearInterval(interval);
                        return;
                    }
                    updateProgress(progress);
                    
                    // تحديث الرسالة بناءً على التقدم
                    if (progress < 20) {
                        updateLoadingMessage('🚀 Starting AI engine...', 'rocket');
                    } else if (progress < 40) {
                        updateLoadingMessage('📚 Loading knowledge base...', 'book');
                    } else if (progress < 60) {
                        updateLoadingMessage('🤖 Generating content...', 'robot');
                    } else if (progress < 80) {
                        updateLoadingMessage('🎨 Designing slides...', 'palette');
                    }
                }, speed);
            }
            
            // ========== 9. إعادة تعيين النموذج ==========
            function resetForm() {
                titleInput.value = '';
                slidesInput.value = 8;
                slideValue.textContent = '8';
                document.querySelector('input[name="prompt_style"][value="prompt1"]').checked = true;
                
                // إعادة تعيين ملف PDF
                uploadedPDF = null;
                pdfFileInput.value = '';
                fileInfo.classList.remove('show');
                
                updateCharCounter();
                
                // إضافة تأثير بسيط
                titleInput.focus();
            }
            
            // ========== 10. نظام الإشعارات ==========
            function showNotification(message, type = 'info') {
                const icons = {
                    success: 'check-circle',
                    error: 'exclamation-circle',
                    info: 'info-circle',
                    warning: 'exclamation-triangle'
                };
                
                const colors = {
                    success: 'linear-gradient(135deg, #c6f6d5 0%, #9ae6b4 100%)',
                    error: 'linear-gradient(135deg, #fed7d7 0%, #feb2b2 100%)',
                    info: 'linear-gradient(135deg, #bee3f8 0%, #90cdf4 100%)',
                    warning: 'linear-gradient(135deg, #feebc8 0%, #fbd38d 100%)'
                };
                
                const textColors = {
                    success: '#22543d',
                    error: '#9b2c2c',
                    info: '#2c5282',
                    warning: '#744210'
                };
                
                // إنشاء الإشعار
                const notification = document.createElement('div');
                notification.className = `notification ${type}`;
                notification.style.cssText = `
                    position: fixed;
                    top: 30px;
                    right: 30px;
                    padding: 20px 25px;
                    border-radius: 12px;
                    display: flex;
                    align-items: center;
                    gap: 15px;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.2);
                    z-index: 10000;
                    animation: slideInRight 0.5s cubic-bezier(0.4, 0, 0.2, 1);
                    max-width: 450px;
                    backdrop-filter: blur(10px);
                    border: 1px solid rgba(255,255,255,0.2);
                    background: ${colors[type] || colors.info};
                    color: ${textColors[type] || textColors.info};
                `;
                
                notification.innerHTML = `
                    <i class="fas fa-${icons[type] || 'info-circle'}" style="font-size: 1.5rem;"></i>
                    <span style="flex: 1; font-weight: 500;">${message}</span>
                    <button class="notification-close" style="background: none; border: none; cursor: pointer; padding: 0; margin-left: 10px;">
                        <i class="fas fa-times" style="font-size: 1.2rem;"></i>
                    </button>
                `;
                
                document.body.appendChild(notification);
                
                // إغلاق يدوي
                const closeBtn = notification.querySelector('.notification-close');
                closeBtn.addEventListener('click', () => {
                    notification.style.animation = 'slideOutRight 0.3s ease forwards';
                    setTimeout(() => notification.remove(), 300);
                });
                
                // إغلاق تلقائي بعد 5 ثوان
                setTimeout(() => {
                    if (notification.parentNode) {
                        notification.style.animation = 'slideOutRight 0.3s ease forwards';
                        setTimeout(() => notification.remove(), 300);
                    }
                }, 5000);
                
                // أصوات الإشعارات (اختياري)
                if (type === 'success') {
                    playSuccessSound();
                } else if (type === 'error') {
                    playErrorSound();
                }
            }
            
            function showSuccess(message) {
                showNotification(message, 'success');
            }
            
            function showError(message) {
                showNotification(message, 'error');
            }
            
            // ========== 11. تأثيرات خاصة ==========
            function addShakeEffect(element) {
                element.style.animation = 'shake 0.5s ease';
                setTimeout(() => {
                    element.style.animation = '';
                }, 500);
            }
            
            function showConfetti() {
                const confettiCount = 50;
                const colors = ['#667eea', '#764ba2', '#ff6b6b', '#48bb78', '#ed8936'];
                
                for (let i = 0; i < confettiCount; i++) {
                    const confetti = document.createElement('div');
                    confetti.className = 'confetti';
                    
                    const size = Math.random() * 10 + 5;
                    const color = colors[Math.floor(Math.random() * colors.length)];
                    const left = Math.random() * 100;
                    const duration = Math.random() * 2 + 2;
                    const delay = Math.random() * 1;
                    
                    confetti.style.cssText = `
                        position: fixed;
                        width: ${size}px;
                        height: ${size}px;
                        background-color: ${color};
                        border-radius: ${Math.random() > 0.5 ? '50%' : '0'};
                        left: ${left}%;
                        top: -20px;
                        z-index: 10001;
                        animation: confetti-fall ${duration}s linear ${delay}s forwards;
                    `;
                    
                    document.body.appendChild(confetti);
                    
                    // إزالة الكونفيتي بعد الانتهاء
                    setTimeout(() => {
                        confetti.remove();
                    }, (duration + delay) * 1000);
                }
            }
            
            // ========== 12. تأثيرات صوتية (اختياري) ==========
            function playSuccessSound() {
                // يمكن إضافة صوت نجاح هنا
                console.log('🎵 Success sound played');
            }
            
            function playErrorSound() {
                // يمكن إضافة صوت خطأ هنا
                console.log('🎵 Error sound played');
            }
            
            // ========== 13. إضافة أنميشن الـ shake ==========
            const style = document.createElement('style');
            style.textContent = `
                @keyframes shake {
                    0%, 100% { transform: translateX(0); }
                    10%, 30%, 50%, 70%, 90% { transform: translateX(-5px); }
                    20%, 40%, 60%, 80% { transform: translateX(5px); }
                }
                
                @keyframes confetti-fall {
                    0% {
                        transform: translateY(-100px) rotate(0deg);
                        opacity: 1;
                    }
                    100% {
                        transform: translateY(100vh) rotate(720deg);
                        opacity: 0;
                    }
                }
            `;
            document.head.appendChild(style);
            
            // ========== 14. تهيئة إضافية ==========
            
            // تحديث عنوان الصفحة ديناميكياً
            const originalTitle = document.title;
            let isTabActive = true;
            
            window.addEventListener('blur', () => {
                isTabActive = false;
            });
            
            window.addEventListener('focus', () => {
                isTabActive = true;
                document.title = originalTitle;
            });
            
            // تحديث عنوان الصفحة عند التحميل
            window.addEventListener('beforeunload', () => {
                if (isLoading) {
                    return "Presentation generation is in progress. Are you sure you want to leave?";
                }
            });
            
            // تهيئة كاملة
            console.log('🚀 AI Presentation Creator Pro initialized successfully!');
            showNotification('Welcome to AI Presentation Creator Pro! 🎉', 'info');
            
            // تأثير كتابة للموقع
            const typewriterText = document.querySelector('.typewriter');
            if (typewriterText) {
                setTimeout(() => {
                    typewriterText.style.animation = 'none';
                    setTimeout(() => {
                        typewriterText.style.animation = 'typing 3.5s steps(40, end), blink-caret .75s step-end infinite';
                    }, 50);
                }, 4000);
            }
            
            // تحديث شريط الشرائح
            slideValue.textContent = slidesInput.value;
        });
        
        // دالة مساعدة لعرض وقت التحميل
        function formatTime(seconds) {
            const mins = Math.floor(seconds / 60);
            const secs = Math.floor(seconds % 60);
            return `${mins}:${secs.toString().padStart(2, '0')}`;
        }
    </script>
</body>
</html>
'''


# ==================== FLASK APP ====================
def create_app():
    app = Flask(__name__)
    app.config['SECRET_KEY'] = 'ai-presentation-creator-secret-key'
    app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32MB max upload

    # Components are created lazily
    content_gen = None
    ppt_creator = None

    def initialize_ai_components():
        """Initialize AI components on demand"""
        nonlocal content_gen, ppt_creator
        try:
            print("🚀 Initializing AI components...")
            content_gen = ContentGenerator(MODEL_PATH)
            ppt_creator = PresentationCreator(content_gen)
            return True
        except Exception as e:
            print(f"❌ Failed to initialize AI components: {e}")
            return False

    @app.after_request
    def cleanup_temp_files_after_request(response):
        try:
            for file_path in list(TEMP_FILES):
                if os.path.exists(file_path):
                    os.remove(file_path)
                    TEMP_FILES.remove(file_path)
                    print(f"🗑️  Cleaned up temp file (after_request): {file_path}")
        except Exception as e:
            print(f"⚠️  Error in after_request cleanup: {e}")
        return response

    @app.route('/')
    def index():
        return render_template_string(HTML_TEMPLATE)

    @app.route('/generate', methods=['POST'])
    def generate():
        try:
            nonlocal content_gen, ppt_creator
            if content_gen is None or ppt_creator is None:
                print("🔄 Initializing AI model on first request...")
                if not initialize_ai_components():
                    return jsonify({'error': 'Failed to initialize AI model. Please check if model file exists.'}), 500

            title = (request.form.get('title') or '').strip()
            slides_str = request.form.get('slides', '8')
            prompt_style = request.form.get('prompt_style', 'prompt1')
            use_pdf_flag = request.form.get('use_pdf', '0') == '1'

            try:
                slides = int(slides_str)
            except ValueError:
                return jsonify({'error': 'Slides must be an integer.'}), 400

            if not title or len(title) < 3:
                return jsonify({'error': 'Please enter a valid title (minimum 3 characters)'}), 400

            if slides < 3 or slides > 12:
                return jsonify({'error': 'Number of slides must be between 3 and 12'}), 400

            # PDF (اختياري + اختياري الاستخدام)
            pdf_file = request.files.get('pdf')
            if use_pdf_flag and pdf_file and pdf_file.filename:
                tmp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
                pdf_path = tmp_pdf.name
                pdf_file.save(pdf_path)
                tmp_pdf.close()
                TEMP_FILES.append(pdf_path)
                pdf_text = extract_text_from_pdf(pdf_path)
                content_gen.set_pdf_from_text(pdf_text)
            else:
                # لم يتم اختيار استخدام PDF أو لم يتم رفع ملف
                content_gen.clear_pdf_index()

            print(f"🎯 Starting generation: '{title}' ({slides} slides, {prompt_style}, use_pdf={use_pdf_flag})")
            start_time = time.time()

            file_path = ppt_creator.create_presentation(title, slides, prompt_style)

            end_time = time.time()
            print(f"⏱️  Generation completed in {end_time - start_time:.2f} seconds")

            return send_file(
                file_path,
                as_attachment=True,
                download_name=f"{title.replace(' ', '_')}.pptx",
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

        except FileNotFoundError:
            return jsonify({'error': f'Model file not found: {MODEL_PATH}. Please download a model and update MODEL_PATH.'}), 500
        except ValueError as e:
            return jsonify({'error': f'Invalid input: {str(e)}'}), 400
        except Exception as e:
            print(f"❌ Unexpected error during generation: {e}")
            return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500

    @app.route('/health')
    def health():
        status = {
            'status': 'ok',
            'model_loaded': content_gen is not None,
            'model_path': MODEL_PATH,
            'model_exists': os.path.exists(MODEL_PATH),
            'temp_files': len(TEMP_FILES)
        }
        return jsonify(status)

    @app.route('/cleanup', methods=['POST'])
    def manual_cleanup():
        """Manually cleanup temp files"""
        count = len(TEMP_FILES)
        cleanup_temp_files()
        return jsonify({'message': f'Cleaned up {count} temporary files'})

    return app


def main():
    print("=" * 60)
    print("🤖 AI PRESENTATION CREATOR (with PDF RAG TF-IDF)")
    print("=" * 60)
    print(f"📂 Model path: {MODEL_PATH}")
    print(f"📁 Model exists: {'✅' if os.path.exists(MODEL_PATH) else '❌'}")

    if not os.path.exists(MODEL_PATH):
        print("\n⚠️  WARNING: Model file not found!")
        print("Please download a GGUF model and:")
        print("1. Place it in the current directory")
        print("2. Update MODEL_PATH in the code")
        print("\nRecommended models:")
        print("- qwen2.5-1.5b-instruct-q4_K_M.gguf (fast)")
        print("- llama-2-7b-chat.Q4_K_M.gguf (balanced)")
        print("- mistral-7b-instruct-v0.2.Q4_K_M.gguf (good quality)")
        print("\nYou can download models from:")
        print("- https://huggingface.co/TheBloke")
        print("- https://huggingface.co/models?sort=trending&search=gguf")

    app = create_app()

    def open_browser():
        time.sleep(1.5)
        try:
            webbrowser.open("http://127.0.0.1:5000")
        except Exception:
            print("⚠️  Could not open browser automatically")

    threading.Thread(target=open_browser, daemon=True).start()

    print("\n" + "=" * 60)
    print("🚀 Starting server...")
    print("🌐 Open your browser to: http://127.0.0.1:5000")
    print("⏳ First request will load the AI model (may take 30-60 seconds)")
    print("🛑 Press Ctrl+C to stop the server")
    print("=" * 60)

    try:
        with simple_server.make_server("127.0.0.1", 5000, app) as httpd:
            httpd.serve_forever()
    except KeyboardInterrupt:
        print("\n\n👋 Server stopped by user")
        print("🧹 Cleaning up temporary files...")
        cleanup_temp_files()
        print("✅ Cleanup complete. Goodbye!")
    except Exception as e:
        print(f"\n❌ Server error: {e}")
        cleanup_temp_files()


if __name__ == "__main__":
    main()
